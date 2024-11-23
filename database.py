import os
import uuid
import subprocess
import numpy as np
import tempfile
from pdf2image import convert_from_path
from pptx import Presentation
from PIL import Image, ImageEnhance, ImageFilter
import pytesseract
import pinecone
from langchain.schema import Document
from langchain_pinecone import PineconeVectorStore
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_upstage import UpstageEmbeddings
from dotenv import load_dotenv
import pandas as pd
from docx import Document as DocxDocument
import fitz  # PyMuPDF
import platform
from logger import logger
import json

# 환경 변수 로드
load_dotenv()

# Pinecone 클라이언트 생성
pc = pinecone.Pinecone(api_key=os.getenv('PINECONE_API_KEY'), environment=os.getenv('PINECONE_ENV'))

# OS에 따른 LibreOffice 경로 설정
if platform.system() == "Windows":
    pytesseract.pytesseract.tesseract_cmd = "C:\\Program Files\\Tesseract-OCR\\tesseract.exe"
    os.environ["TESSDATA_PREFIX"] = "C:\\Program Files\\Tesseract-OCR\\tessdata"
    poppler_path = "C:\\Users\\uracle\\Desktop\\python-workspace\\ai\\rag-atoz\\libs\\poppler\\Library\\bin"
    LIBREOFFICE_PATH = "C:\\Program Files\\LibreOffice\\program\\soffice.exe"
else:
    pytesseract.pytesseract.tesseract_cmd = "/usr/bin/tesseract"
    os.environ["TESSDATA_PREFIX"] = "/usr/share/tesseract-ocr/5/tessdata"
    poppler_path = "/usr/bin"
    LIBREOFFICE_PATH = "/usr/lib/libreoffice/program/soffice"  # linux 경로
    

# 파일 형식별 처리 함수 정의

def process_text_file(file_path):
    """TXT 파일의 내용을 읽어 문자열로 반환합니다."""
    with open(file_path, "r", encoding="utf-8") as file:
        text = file.read()
    print(f"TXT 파일 내용:\n{text}\n{'-'*50}")
    return text

def process_json_file(file_path):
    """JSON 파일의 내용을 읽어 문자열로 반환합니다."""
    with open(file_path, "r", encoding="utf-8") as file:
        datalist = json.load(file)
    print(f"JSON 파일 내용:\n{file}\n{'-'*50}")
    return datalist

def process_docx_file(file_path):
    """DOCX 파일의 각 단락을 읽어 문자열로 결합하여 반환하며, 파일명을 metadata에 추가합니다."""
    doc = DocxDocument(file_path)
    document_text = "\n".join([para.text for para in doc.paragraphs])
    return document_text  # 파일명은 반환하지 않고 호출 함수에서 전달

def process_doc_file(file_path):
    """DOC 파일을 LibreOffice를 사용하여 DOCX로 변환한 후, 텍스트를 추출하여 반환합니다."""
    # DOC 파일을 DOCX로 변환
    docx_path = file_path.replace(".doc", ".docx")
    result = subprocess.run(
        [LIBREOFFICE_PATH, "--headless", "--convert-to", "docx", file_path, "--outdir", os.path.dirname(file_path)],
        capture_output=True, text=True
    )
    
    # 변환 실패 시 예외 발생
    if result.returncode != 0:
        print("LibreOffice DOC to DOCX 변환 실패:", result.stderr)
        raise Exception("DOC to DOCX 변환 실패")
    
    # 변환된 DOCX 파일에서 텍스트 추출
    document_text = process_docx_file(docx_path)
    print(f"DOC 파일 내용:\n{document_text}\n{'-'*50}")
    
    # 변환된 DOCX 파일 삭제
    os.remove(docx_path)
    
    return document_text

def process_xlsx_file(file_path):
    """엑셀 파일의 각 시트를 개별 Document로 생성하여 반환하고, 긴 텍스트는 분할하여 저장합니다."""
    sheet_documents = []
    xls = pd.ExcelFile(file_path)

    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name)
        sheet_text = " ".join(df.astype(str).values.flatten())
        
        # 텍스트 분할: 1000 토큰으로 청크를 분할, 200 토큰 중첩
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        split_texts = text_splitter.split_text(sheet_text)
        
        # 각 분할된 텍스트를 Document로 생성
        for text in split_texts:
            sheet_documents.append(Document(page_content=text, metadata={"sheet_name": sheet_name}))

    return sheet_documents

def process_pdf_file(file_path):
    """PDF 파일의 모든 페이지에서 텍스트를 추출하여 하나의 문자열로 반환합니다."""
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    print(f"PDF 파일 내용:\n{text}\n{'-'*50}")
    return text

def process_image_file(file_path):
    """이미지 파일에서 OCR을 사용해 텍스트를 추출하여 문자열로 반환합니다."""
    image = Image.open(file_path)
    image = image.convert("L")  # 흑백 변환
    image = ImageEnhance.Contrast(image).enhance(2)  # 대비 증가
    image = image.filter(ImageFilter.SHARPEN)  # 샤프닝 필터 적용
    text = pytesseract.image_to_string(image, lang="kor").strip()
    print(f"이미지 파일 OCR 내용:\n{text}\n{'-'*50}")
    return text

def preprocess_image_for_ocr(image_path):
    """이미지 전처리: 대비 증가, 흑백 변환, 샤프닝 적용"""
    image = Image.open(image_path)
    image = image.convert("L")  # 흑백 변환
    image = ImageEnhance.Contrast(image).enhance(2)  # 대비 증가
    image = image.filter(ImageFilter.SHARPEN)  # 샤프닝 필터 적용
    return image

def convert_pptx_to_pdf_images(pptx_path, save_dir="C:\\Users\\uracle\\Desktop\\python-workspace\\ai\\rag-atoz\\images\\poppler"):
    """PPTX 파일을 PDF로 변환 후, PDF의 각 페이지를 이미지로 변환하여 상대 경로로 이미지 파일 경로를 반환합니다."""
    if platform.system() == "Windows":
        save_dir = "C:\\Users\\uracle\\Desktop\\python-workspace\\ai\\rag-atoz\\images\\poppler"
    else:
        save_dir = "/ai/rag-atoz/images/poppler"
    os.makedirs(save_dir, exist_ok=True)
    
    # PDF 파일 경로 설정 (temp 폴더에 임시로 저장)
    pdf_path = os.path.join("temp", f"{os.path.splitext(os.path.basename(pptx_path))[0]}.pdf")
    
    # LibreOffice를 사용하여 PPTX를 PDF로 변환
    result = subprocess.run(
        [LIBREOFFICE_PATH, "--headless", "--convert-to", "pdf", pptx_path, "--outdir", os.path.dirname(pdf_path)],
        capture_output=True, text=True
    )
    
    if result.returncode != 0:
        logger.log_custom("LibreOffice 변환 실패:%s", result.stderr)
        raise Exception("PPTX to PDF 변환 실패")

    if not os.path.exists(pdf_path):
        raise Exception("PDF 파일이 생성되지 않았습니다.")
    
    images = convert_from_path(pdf_path, poppler_path=poppler_path)
    image_paths = []

    for i, img in enumerate(images):
        # 상대 경로로 이미지 저장
        image_name = f"{os.path.splitext(os.path.basename(pptx_path))[0]}_page_{i + 1}.png"
        image_path = os.path.join(save_dir, image_name)
        img.save(image_path, "PNG")
        
        # 이미지 상대 경로로 리스트에 추가
        image_paths.append(image_path.replace("\\", "/"))  # 경로에 슬래시(`/`) 사용

    os.remove(pdf_path)  # 임시 PDF 파일 삭제
    return image_paths

def deptJson(datalist):
    documents_to_upload = []
    vectors_to_upload = []
    for entry in datalist:
        data = entry 
        vector_id = str(uuid.uuid4())
        
        # 랜덤 3차원 임베딩 값 생성
        vector_values = np.random.rand(3).tolist()
    
        # Document 객체 생성
        docu = Document(
            page_content=f"""
            -{data['이름']} 사원정보:
            {data["이름"]}님의 소속은 {data["소속"]} 입니다. 사원번호는 {data["사원번호"]} 이며, 
            부서는 {data["부서명"]} 입니다. 직위(직급)은 {data["직위(직급)"]} 이며, 
            직책은 {data["직책"]} 입니다. 휴대전화는 {data["연락처(휴대전화)"]} 이며,
            이메일은 {data["이메일주소"]} 입니다.
            """, 
            metadata={
                "이름": data["이름"],
                "회사명": data["회사명"],
                "부서명": data["부서명"],
                "직위(직급)": data["직위(직급)"],
                "직책": data["직책"],
                "사원번호": data["사원번호"],
                "연락처(휴대전화)": data["연락처(휴대전화)"],
                "소속": data["소속"],
                "이메일주소": data["이메일주소"],
                "생년월일": data["생년월일"],
                "내선전화번호": data["내선전화번호"],
                "팩스번호": data["팩스번호"],
                "source": "../dept-user-markdown-table.json",
                "category": "사원정보",
                "table_title": "사원정보 표",
                "type": "markdown_table",
                "updated": "2024.10.29",
                "description": """이 표는 유라클, (주)지네트웍스, 에이네트웍스 등 3개 회사의 사원정보 표 입니다.
                이름, 회사명, 부서명, 직위(직급), 직책, 사원번호, 연락처, 소속, 이메일, 내선번호 등의 정보가 있습니다.
                """
            }
        )
        documents_to_upload.append(docu)  # Document 리스트에 추가
    return documents_to_upload
    

def process_and_store_document(uploaded_file, file_type):
    """파일을 읽어 Pinecone에 저장합니다. 파일 형식에 따라 적절한 처리 함수를 호출합니다."""
    logger.log_custom("process_and_store_document 함수 호출")
    try:
        
        # 원본 파일명을 저장
        file_name = uploaded_file.name
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as temp_file:
            temp_file.write(uploaded_file.read())
            temp_file_path = temp_file.name

        # Document 객체를 저장할 리스트
        documents = []

        if file_type == "pptx":
            # PPTX 파일의 슬라이드를 PDF로 변환 후 각 페이지를 이미지로 변환
            image_files = convert_pptx_to_pdf_images(temp_file_path)  # PDF -> 이미지 변환
            logger.log_custom("temp_file_path:%s",temp_file_path)
            ppt = Presentation(temp_file_path)
            logger.log_custom("Presentation")
            for i, slide in enumerate(ppt.slides):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_text.append(shape.text)
                 
                # OCR 수행
                ocr_text_content = ""
                if i < len(image_files):
                    preprocessed_image = preprocess_image_for_ocr(image_files[i])
                    logger.log_custom("preprocess_image_for_ocr:%s번째",i)
                    ocr_text_content = pytesseract.image_to_string(preprocessed_image, lang="kor").strip()
                    logger.log_custom("ocr_text_content")

                # 슬라이드 텍스트와 OCR 텍스트 결합하여 Document 객체 생성
                combined_content = f"{' '.join(slide_text)}\n{ocr_text_content}"
                logger.log_custom(f"PPTX 파일 슬라이드 {i+1} 내용:\n{combined_content}\n{'-'*50}") 

                

                # 파일명과 확장자 추출
                image_file_name, image_file_extension = os.path.splitext(os.path.basename(image_files[i]))
                image_url = os.getenv('IMAGE_URL')+"/poppler/"+image_file_name+image_file_extension
                documents.append(Document(
                    page_content=combined_content+f"\n\n![Image]({image_url}) \"참고 이미지\"" ,
                    metadata={
                        "slide_index": i + 1,
                        #"image_path": image_files[i],
                        "image_path": image_url,
                        "source": file_name
                    }
                ))

        elif file_type == "xlsx":
            # 엑셀 파일 각 시트를 Document로 생성하고, 파일명 metadata에 추가
            sheet_documents = process_xlsx_file(temp_file_path)
            documents = [Document(page_content=sheet.page_content, metadata={"source": file_name, "sheet_name": sheet.metadata["sheet_name"]}) for sheet in sheet_documents]
        
        elif file_type == "docx":
            document_text = process_docx_file(temp_file_path)
            print(f"DOCX 파일 내용:\n{document_text}\n{'-'*50}")
            
            # 텍스트 분할
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            split_texts = text_splitter.split_text(document_text)
             # 분할된 텍스트 각각에 Document 생성 및 metadata에 원본 파일명을 포함한 source 추가
            documents = [
                Document(
                    page_content=text,
                    metadata={
                        "source": file_name
                    }
                ) for index, text in enumerate(split_texts)
            ]

        elif file_type == "doc":
            document_text = process_doc_file(temp_file_path)
            print(f"DOC 파일 내용:\n{document_text}\n{'-'*50}")
            
            # 텍스트 분할
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            split_texts = text_splitter.split_text(document_text)
            documents = [Document(page_content=text, metadata={"source": file_name}) for text in split_texts]

        elif file_type == "pdf":
            document_text = process_pdf_file(temp_file_path)
            print(f"PDF 파일 내용:\n{document_text}\n{'-'*50}")
            
            # 텍스트 분할
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            split_texts = text_splitter.split_text(document_text)
            documents = [Document(page_content=text, metadata={"source": file_name}) for text in split_texts]

        elif file_type in ["png", "jpg", "jpeg"]:
            document_text = process_image_file(temp_file_path)
            print(f"이미지 파일 OCR 내용:\n{document_text}\n{'-'*50}")
            
            # 이미지 OCR 결과를 Document 객체로 생성
            documents = [Document(page_content=document_text, metadata={"source": file_name})]

        elif file_type == "txt":
            document_text =  process_text_file(temp_file_path)
            print(f"텍스트 파일 내용:\n{document_text}\n{'-'*50}")
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            split_texts = text_splitter.split_text(document_text)
            documents = [Document(page_content=text, metadata={"source": file_name}) for text in split_texts]

        elif file_type == "json":
            if file_name == 'dept-user-markdown-table.json':
                datalist = process_json_file(temp_file_path)
                documents = deptJson(datalist)
            else:
                datalist = process_json_file(temp_file_path)
                 
                # JSON 데이터가 리스트 형태일 경우 처리
                if isinstance(datalist, list):
                    documents = []
                    for item in datalist:
                        # JSON의 각 항목(item)을 텍스트와 메타데이터로 변환
                        page_content = json.dumps(item, ensure_ascii=False, indent=4)
                        metadata = {"source": file_name}

                        # JSON 내부에 특정 키를 메타데이터로 추가 가능
                        if isinstance(item, dict):
                            metadata.update({key: item[key] for key in item.keys() if key != "content"})

                        # Document 생성
                        documents.append(Document(page_content=page_content, metadata=metadata))
                
                # JSON 데이터가 딕셔너리 형태일 경우 처리
                elif isinstance(datalist, dict):
                    page_content = json.dumps(datalist, ensure_ascii=False, indent=4)
                    metadata = {"source": file_name}
                    metadata.update(datalist)  # JSON의 키-값을 메타데이터로 추가
                    documents = [Document(page_content=page_content, metadata=metadata)]

                else:
                    raise ValueError("지원하지 않는 JSON 데이터 형식입니다.")

        # Pinecone에 저장
        embeddings = UpstageEmbeddings(model="solar-embedding-1-large")
        index_name = os.getenv('INDEX_NAME')
        
        # 인덱스가 존재하는지 확인
        if index_name not in pc.list_indexes().names():
            pc.create_index(index_name, dimension=embeddings.dimension)
        
        if len(documents) > 0:
            logger.log_custom("DB INSERT LENGTH:%s",len(documents))
            logger.log_custom("DB LOAD") 
            database = PineconeVectorStore.from_documents(documents, embeddings, index_name=index_name)
        
        logger.log_custom("임베딩이 성공적으로 저장되었습니다!")
        return True

    except Exception as e:
        logger.log_custom("임베딩 저장 실패:%s",str(e))
        return False